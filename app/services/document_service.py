"""
DocMind AI — Document parsing service
Extracts plain text from PDF, DOCX, PPTX, and TXT files.
"""

from __future__ import annotations

import os
from pathlib import Path


def extract_text(file_path: str) -> tuple[str, int]:
    """
    Extract plain text from a document.

    Returns:
        (text, page_count)
    """
    ext = Path(file_path).suffix.lower().lstrip(".")
    if ext == "pdf":
        return _extract_pdf(file_path)
    elif ext == "docx":
        return _extract_docx(file_path), 1
    elif ext == "pptx":
        return _extract_pptx(file_path), 1
    elif ext == "txt":
        return _extract_txt(file_path), 1
    else:
        return "", 0


def _extract_pdf(file_path: str) -> tuple[str, int]:
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        pages = [page.get_text() for page in doc]
        return "\n".join(pages), len(pages)
    except Exception as exc:  # noqa: BLE001
        return f"[PDF extraction error: {exc}]", 0


def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as exc:  # noqa: BLE001
        return f"[DOCX extraction error: {exc}]"


def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text)
        return "\n".join(lines)
    except Exception as exc:  # noqa: BLE001
        return f"[PPTX extraction error: {exc}]"


def _extract_txt(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read()
    except Exception as exc:  # noqa: BLE001
        return f"[TXT extraction error: {exc}]"


def count_words(text: str) -> int:
    return len(text.split())


def estimate_reading_time(word_count: int, wpm: int = 200) -> int:
    """Return estimated reading time in minutes (minimum 1)."""
    return max(1, round(word_count / wpm))
